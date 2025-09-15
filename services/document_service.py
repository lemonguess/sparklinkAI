"""文档处理服务"""
import os
import logging
from typing import List, Dict, Any, Optional
import mimetypes
from pathlib import Path

# 文档处理相关库
try:
    import PyPDF2
    from docx import Document as DocxDocument
    from pptx import Presentation
    from PIL import Image
    import pytesseract
except ImportError:
    # 如果没有安装这些库，会在运行时提示
    pass

from core.config import settings
from models.database import Document

logger = logging.getLogger(__name__)

class DocumentService:
    """文档处理服务类"""
    
    def __init__(self):
        self.supported_types = {
            'application/pdf': self._extract_pdf_content,
            'application/msword': self._extract_doc_content,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._extract_docx_content,
            'application/vnd.ms-powerpoint': self._extract_ppt_content,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._extract_pptx_content,
            'text/plain': self._extract_text_content,
            'text/markdown': self._extract_text_content,
            'image/jpeg': self._extract_image_content,
            'image/png': self._extract_image_content,
            'image/gif': self._extract_image_content
        }
    
    def extract_content(self, file_path: str, file_type: str) -> str:
        """提取文档内容"""
        try:
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"文件不存在: {file_path}")
            
            if file_type not in self.supported_types:
                raise ValueError(f"不支持的文件类型: {file_type}")
            
            # 调用对应的提取方法
            extract_method = self.supported_types[file_type]
            content = extract_method(file_path)
            
            if not content or not content.strip():
                raise ValueError("文档内容为空")
            
            logger.info(f"成功提取文档内容: {file_path}, 长度: {len(content)}")
            return content.strip()
            
        except Exception as e:
            logger.error(f"提取文档内容失败: {file_path}, 错误: {e}")
            raise
    
    def _extract_pdf_content(self, file_path: str) -> str:
        """提取PDF内容"""
        try:
            content = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    content += page.extract_text() + "\n"
            return content
        except Exception as e:
            logger.error(f"PDF内容提取失败: {e}")
            # 如果PyPDF2失败，可以尝试其他方法或OCR
            return self._extract_pdf_with_ocr(file_path)
    
    def _extract_pdf_with_ocr(self, file_path: str) -> str:
        """使用OCR提取PDF内容（备用方案）"""
        try:
            # 这里可以集成TextIn OCR API或其他OCR服务
            # 暂时返回简单提示
            return "PDF内容提取需要OCR服务支持，请配置TextIn API密钥。"
        except Exception as e:
            logger.error(f"OCR提取PDF失败: {e}")
            raise
    
    def _extract_docx_content(self, file_path: str) -> str:
        """提取DOCX内容"""
        try:
            doc = DocxDocument(file_path)
            content = ""
            for paragraph in doc.paragraphs:
                content += paragraph.text + "\n"
            
            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = "\t".join([cell.text for cell in row.cells])
                    content += row_text + "\n"
            
            return content
        except Exception as e:
            logger.error(f"DOCX内容提取失败: {e}")
            raise
    
    def _extract_doc_content(self, file_path: str) -> str:
        """提取DOC内容（旧版Word文档）"""
        try:
            # DOC格式比较复杂，这里简化处理
            # 实际项目中可能需要使用python-docx2txt或其他库
            return "DOC格式文档需要额外的处理库支持。"
        except Exception as e:
            logger.error(f"DOC内容提取失败: {e}")
            raise
    
    def _extract_pptx_content(self, file_path: str) -> str:
        """提取PPTX内容"""
        try:
            prs = Presentation(file_path)
            content = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                content += f"\n=== 幻灯片 {slide_num} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            
            return content
        except Exception as e:
            logger.error(f"PPTX内容提取失败: {e}")
            raise
    
    def _extract_ppt_content(self, file_path: str) -> str:
        """提取PPT内容（旧版PowerPoint）"""
        try:
            # PPT格式需要特殊处理
            return "PPT格式文档需要额外的处理库支持。"
        except Exception as e:
            logger.error(f"PPT内容提取失败: {e}")
            raise
    
    def _extract_text_content(self, file_path: str) -> str:
        """提取纯文本内容"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                with open(file_path, 'r', encoding='gbk') as file:
                    return file.read()
            except Exception:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
        except Exception as e:
            logger.error(f"文本内容提取失败: {e}")
            raise
    
    def _extract_image_content(self, file_path: str) -> str:
        """提取图片内容（OCR）"""
        try:
            # 使用TextIn OCR API或本地OCR
            if settings.TEXTIN_API_KEY:
                return self._extract_image_with_textin(file_path)
            else:
                return self._extract_image_with_tesseract(file_path)
        except Exception as e:
            logger.error(f"图片内容提取失败: {e}")
            raise
    
    def _extract_image_with_textin(self, file_path: str) -> str:
        """使用TextIn API提取图片文字"""
        try:
            # 这里应该调用TextIn OCR API
            # 暂时返回占位符
            return "图片OCR功能需要配置TextIn API密钥。"
        except Exception as e:
            logger.error(f"TextIn OCR失败: {e}")
            raise
    
    def _extract_image_with_tesseract(self, file_path: str) -> str:
        """使用Tesseract OCR提取图片文字"""
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            return text
        except Exception as e:
            logger.error(f"Tesseract OCR失败: {e}")
            return "本地OCR功能需要安装Tesseract。"
    
    def split_content(
        self,
        content: str,
        chunk_size: int = 512,
        chunk_overlap: int = 50
    ) -> List[str]:
        """分割文档内容"""
        try:
            if not content or not content.strip():
                return []
            
            # 简单的分块策略：按段落和句子分割
            chunks = []
            
            # 首先按段落分割
            paragraphs = content.split('\n\n')
            
            current_chunk = ""
            
            for paragraph in paragraphs:
                paragraph = paragraph.strip()
                if not paragraph:
                    continue
                
                # 如果当前段落加上现有chunk超过大小限制
                if len(current_chunk) + len(paragraph) > chunk_size:
                    if current_chunk:
                        chunks.append(current_chunk.strip())
                        # 保留重叠部分
                        if chunk_overlap > 0 and len(current_chunk) > chunk_overlap:
                            current_chunk = current_chunk[-chunk_overlap:] + "\n" + paragraph
                        else:
                            current_chunk = paragraph
                    else:
                        # 单个段落就超过限制，需要进一步分割
                        sub_chunks = self._split_long_paragraph(paragraph, chunk_size, chunk_overlap)
                        chunks.extend(sub_chunks[:-1])  # 除了最后一个
                        current_chunk = sub_chunks[-1] if sub_chunks else ""
                else:
                    if current_chunk:
                        current_chunk += "\n\n" + paragraph
                    else:
                        current_chunk = paragraph
            
            # 添加最后一个chunk
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
            
            # 过滤掉太短的chunks
            chunks = [chunk for chunk in chunks if len(chunk.strip()) > 20]
            
            logger.info(f"文档分块完成: 原长度 {len(content)}, 分成 {len(chunks)} 块")
            return chunks
            
        except Exception as e:
            logger.error(f"文档分块失败: {e}")
            raise
    
    def _split_long_paragraph(
        self,
        paragraph: str,
        chunk_size: int,
        chunk_overlap: int
    ) -> List[str]:
        """分割长段落"""
        chunks = []
        
        # 按句子分割
        sentences = self._split_sentences(paragraph)
        
        current_chunk = ""
        
        for sentence in sentences:
            if len(current_chunk) + len(sentence) > chunk_size:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                    # 保留重叠
                    if chunk_overlap > 0 and len(current_chunk) > chunk_overlap:
                        current_chunk = current_chunk[-chunk_overlap:] + sentence
                    else:
                        current_chunk = sentence
                else:
                    # 单个句子就超过限制，强制分割
                    if len(sentence) > chunk_size:
                        for i in range(0, len(sentence), chunk_size - chunk_overlap):
                            chunk = sentence[i:i + chunk_size]
                            if chunk.strip():
                                chunks.append(chunk.strip())
                    else:
                        current_chunk = sentence
            else:
                if current_chunk:
                    current_chunk += sentence
                else:
                    current_chunk = sentence
        
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def _split_sentences(self, text: str) -> List[str]:
        """分割句子"""
        # 简单的句子分割，基于标点符号
        import re
        
        # 中英文句子结束标点
        sentence_endings = r'[.!?。！？；;]\s*'
        sentences = re.split(sentence_endings, text)
        
        # 过滤空句子
        sentences = [s.strip() for s in sentences if s.strip()]
        
        return sentences
    
    async def process_document(self, document_id: int):
        """处理文档（异步调用Celery任务）"""
        try:
            from services.tasks.document_tasks import process_document_task
            
            # 提交异步任务
            task = process_document_task.delay(document_id)
            
            logger.info(f"文档处理任务已提交: {document_id}, 任务ID: {task.id}")
            
            return {
                "task_id": task.id,
                "document_id": document_id,
                "status": "submitted"
            }
            
        except Exception as e:
            logger.error(f"提交文档处理任务失败: {e}")
            raise
    
    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """获取文件信息"""
        try:
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"文件不存在: {file_path}")
            
            file_stat = os.stat(file_path)
            file_type, _ = mimetypes.guess_type(file_path)
            
            return {
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "file_size": file_stat.st_size,
                "file_type": file_type,
                "created_time": file_stat.st_ctime,
                "modified_time": file_stat.st_mtime,
                "is_supported": file_type in self.supported_types
            }
            
        except Exception as e:
            logger.error(f"获取文件信息失败: {e}")
            raise
    
    def validate_file(self, file_path: str, max_size: int = 10 * 1024 * 1024) -> bool:
        """验证文件"""
        try:
            file_info = self.get_file_info(file_path)
            
            # 检查文件大小
            if file_info["file_size"] > max_size:
                raise ValueError(f"文件大小超过限制: {file_info['file_size']} > {max_size}")
            
            # 检查文件类型
            if not file_info["is_supported"]:
                raise ValueError(f"不支持的文件类型: {file_info['file_type']}")
            
            return True
            
        except Exception as e:
            logger.error(f"文件验证失败: {e}")
            raise